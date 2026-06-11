from __future__ import annotations

import json
import os
import shutil
import tempfile
import unittest

import pytest

# 这些解析/造测试数据用的第三方库都较重且可选 (document_parser 本体也是 optional import,
# 见 core/impl.py)。缺任一就整模块 **skip** 而非 collection error — 让全量 pytest 不被中断;
# 装了则照常跑。
pd = pytest.importorskip("pandas")
xmltodict = pytest.importorskip("xmltodict")
pytest.importorskip("docx")
pytest.importorskip("pptx")
pytest.importorskip("reportlab")

from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

from document_parser_module.core.impl import LocalDocumentParser


class TestLocalDocumentParser(unittest.TestCase):
    """文档解析模块具体实现类的单元测试，覆盖核心功能与异常场景。"""

    def setUp(self):
        self.parser = LocalDocumentParser()
        self.tmp_dir = tempfile.mkdtemp(prefix="doc_parser_test_")

        self.test_txt_path = os.path.join(self.tmp_dir, "test.txt")
        self.test_pdf_path = os.path.join(self.tmp_dir, "test.pdf")
        self.test_docx_path = os.path.join(self.tmp_dir, "test.docx")
        self.test_md_path = os.path.join(self.tmp_dir, "test.md")
        self.test_py_path = os.path.join(self.tmp_dir, "test.py")
        self.test_excel_path = os.path.join(self.tmp_dir, "test.xlsx")
        self.test_ppt_path = os.path.join(self.tmp_dir, "test.pptx")
        self.test_csv_path = os.path.join(self.tmp_dir, "test.csv")
        self.test_json_path = os.path.join(self.tmp_dir, "test.json")
        self.test_xml_path = os.path.join(self.tmp_dir, "test.xml")
        self.test_html_path = os.path.join(self.tmp_dir, "test.html")

        # txt
        with open(self.test_txt_path, "w", encoding="utf-8") as f:
            f.write("hello txt\nline2")

        # pdf (reportlab)
        c = canvas.Canvas(self.test_pdf_path, pagesize=letter)
        c.drawString(100, 750, "hello pdf")
        c.showPage()
        c.drawString(100, 750, "page2")
        c.save()

        # docx
        doc = Document()
        doc.add_paragraph("hello docx")
        doc.add_paragraph("line2")
        doc.save(self.test_docx_path)

        # md
        with open(self.test_md_path, "w", encoding="utf-8") as f:
            f.write("# Title\n\n- a\n- b\n\n**bold**")

        # py
        with open(self.test_py_path, "w", encoding="utf-8") as f:
            f.write("# comment\nprint('hi')\n")

        # excel
        df1 = pd.DataFrame({"a": ["1", "2"], "b": ["x", "y"]})
        df2 = pd.DataFrame({"c": ["3"], "d": ["z"]})
        with pd.ExcelWriter(self.test_excel_path) as writer:
            df1.to_excel(writer, sheet_name="s1", index=False)
            df2.to_excel(writer, sheet_name="s2", index=False)

        # pptx
        pres = Presentation()
        slide = pres.slides.add_slide(pres.slide_layouts[0])
        slide.shapes.title.text = "hello ppt"
        pres.slides.add_slide(pres.slide_layouts[5])
        pres.save(self.test_ppt_path)

        # csv
        df1.to_csv(self.test_csv_path, index=False)

        # json
        with open(self.test_json_path, "w", encoding="utf-8") as f:
            json.dump({"k": "v", "n": 1}, f, ensure_ascii=False)

        # xml
        xml_obj = {"root": {"item": [{"@id": "1", "#text": "v"}]}}
        xml_text = xmltodict.unparse(xml_obj, pretty=True)
        with open(self.test_xml_path, "w", encoding="utf-8") as f:
            f.write(xml_text)

        # html (script/style 内容须被剔除)
        with open(self.test_html_path, "w", encoding="utf-8") as f:
            f.write(
                "<html><head><title>页面标题</title>"
                "<style>body { color: red; }</style>"
                "<script>var secret = 'JS_NOISE';</script></head>"
                "<body><h1>hello html</h1><p>正文段落</p></body></html>"
            )

    def tearDown(self):
        shutil.rmtree(self.tmp_dir, ignore_errors=True)

    def _assert_standard_structure(self, res: dict, ext: str):
        self.assertIn("content", res)
        self.assertIn("file_name", res)
        self.assertIn("meta", res)
        self.assertEqual(res["meta"].get("ext"), ext)
        self.assertTrue(isinstance(res["content"], str))
        self.assertTrue(len(res["content"]) > 0)

    def test_parse_txt_file(self):
        res = self.parser.parse_file(self.test_txt_path)
        self._assert_standard_structure(res, ".txt")
        self.assertIn("hello txt", res["content"])

    def test_parse_pdf_file(self):
        res = self.parser.parse_file(self.test_pdf_path)
        self._assert_standard_structure(res, ".pdf")
        self.assertIn("hello pdf", res["content"])
        self.assertIn("page2", res["content"])

    def test_parse_docx_file(self):
        res = self.parser.parse_file(self.test_docx_path)
        self._assert_standard_structure(res, ".docx")
        self.assertIn("hello docx", res["content"])

    def test_parse_md_file(self):
        res = self.parser.parse_file(self.test_md_path)
        self._assert_standard_structure(res, ".md")
        self.assertIn("Title", res["content"])
        self.assertIn("bold", res["content"])

    def test_parse_html_file(self):
        res = self.parser.parse_file(self.test_html_path)
        self._assert_standard_structure(res, ".html")
        self.assertIn("hello html", res["content"])
        self.assertIn("正文段落", res["content"])
        # script/style 内容不是正文
        self.assertNotIn("JS_NOISE", res["content"])
        self.assertNotIn("color: red", res["content"])

    def test_parse_zip_archive(self):
        # 压缩包是容器: 文本/已知格式成员解析正文, 二进制/嵌套包只列清单
        import zipfile
        zpath = os.path.join(self.tmp_dir, "bundle.zip")
        with zipfile.ZipFile(zpath, "w") as zf:
            zf.writestr("docs/读我.txt", "压缩包成员正文: 火星车")
            zf.writestr("notes.md", "# 标题\n第二成员内容")
            zf.writestr("img.png", b"\x89PNG\x00\x00binaryjunk")
            zf.writestr("inner.zip", b"PK\x03\x04\x00\x00")
        res = self.parser.parse_file(zpath)
        self._assert_standard_structure(res, ".zip")
        c = res["content"]
        self.assertIn("火星车", c)             # 文本成员正文
        self.assertIn("第二成员内容", c)        # md 成员经 parse_file 解析
        self.assertIn("读我.txt", c)           # 清单含中文文件名
        self.assertIn("img.png", c)            # 二进制成员列清单
        self.assertIn("二进制", c)             # 且标注原因
        self.assertIn("嵌套压缩包", c)          # 嵌套包不展开
        self.assertNotIn("\x00", c)            # 正文无二进制泄漏

    def test_parse_nested_zip_one_level(self):
        # 直接嵌套的 zip 展开 1 层; 二层嵌套只列清单不展开
        import io
        import zipfile
        inner2 = io.BytesIO()
        with zipfile.ZipFile(inner2, "w") as z2:
            z2.writestr("deep.txt", "三层深处: 不该出现")
        inner1 = io.BytesIO()
        with zipfile.ZipFile(inner1, "w") as z1:
            z1.writestr("src/main.py", "# 嵌套一层的源码: 蜂鸟引擎")
            z1.writestr("inner2.zip", inner2.getvalue())
        zpath = os.path.join(self.tmp_dir, "outer.zip")
        with zipfile.ZipFile(zpath, "w") as zf:
            zf.writestr("readme.txt", "外层说明")
            zf.writestr("修改版.zip", inner1.getvalue())
        res = self.parser.parse_file(zpath)
        c = res["content"]
        self.assertIn("外层说明", c)
        self.assertIn("蜂鸟引擎", c)          # 一层嵌套包内的文本成员被提取
        self.assertIn("src/main.py", c)       # 嵌套包清单也在
        self.assertNotIn("不该出现", c)        # 二层嵌套不展开
        self.assertIn("层级或数量超限", c)      # 且有明确标注

    def test_parse_tar_gz_archive(self):
        import io
        import tarfile as _tarfile
        tpath = os.path.join(self.tmp_dir, "logs.tar.gz")
        with _tarfile.open(tpath, "w:gz") as tf:
            data = "tar 成员关键词: 风暴眼".encode("utf-8")
            ti = _tarfile.TarInfo("app.log")
            ti.size = len(data)
            tf.addfile(ti, io.BytesIO(data))
        res = self.parser.parse_file(tpath)
        self.assertIn("风暴眼", res["content"])
        self.assertIn("app.log", res["content"])

    def test_parse_unknown_ext_text_fallback(self):
        # 未知后缀 (白名单外) 的文本文件 → 内容嗅探兜底为纯文本, 不再被拒
        log_path = os.path.join(self.tmp_dir, "app.log")
        with open(log_path, "w", encoding="utf-8") as f:
            f.write("2026-06-11 INFO 服务启动\n2026-06-11 ERROR 连接超时")
        res = self.parser.parse_file(log_path)
        self._assert_standard_structure(res, ".log")
        self.assertIn("连接超时", res["content"])


    def test_parse_py_file(self):
        res = self.parser.parse_file(self.test_py_path)
        self._assert_standard_structure(res, ".py")
        self.assertIn("print('hi')", res["content"])

    def test_parse_excel_file(self):
        res = self.parser.parse_file(self.test_excel_path)
        self._assert_standard_structure(res, ".xlsx")
        self.assertIn("[Sheet: s1]", res["content"])
        self.assertIn("[Sheet: s2]", res["content"])

    def test_parse_ppt_file(self):
        res = self.parser.parse_file(self.test_ppt_path)
        self._assert_standard_structure(res, ".pptx")
        self.assertIn("[Slide 1]", res["content"])
        self.assertIn("hello ppt", res["content"])

    def test_parse_csv_file(self):
        res = self.parser.parse_file(self.test_csv_path)
        self._assert_standard_structure(res, ".csv")
        self.assertIn("a\tb", res["content"])

    def test_parse_json_file(self):
        res = self.parser.parse_file(self.test_json_path)
        self._assert_standard_structure(res, ".json")
        self.assertIn('"k": "v"', res["content"])

    def test_parse_xml_file(self):
        res = self.parser.parse_file(self.test_xml_path)
        self._assert_standard_structure(res, ".xml")
        self.assertIn('"root"', res["content"])

    def test_parse_invalid_file_type(self):
        # 内容嗅探后, "拒收"的判据从扩展名清单变成二进制内容 (头部含 NUL)
        invalid_path = os.path.join(self.tmp_dir, "bad.bin")
        with open(invalid_path, "wb") as f:
            f.write(b"MZ\x00\x03\x00\x04")
        with self.assertRaises(Exception) as ctx:
            self.parser.parse_file(invalid_path)
        # RAGException 的 str() 是 message 不含 code; 错误码在 .code 属性上
        self.assertEqual(ctx.exception.code, "UNSUPPORTED_FILE_TYPE")

    def test_parse_folder(self):
        res_list = self.parser.parse_folder(self.tmp_dir)
        # 10 supported files created in setUp
        self.assertGreaterEqual(len(res_list), 10)
        for r in res_list:
            self.assertIn("content", r)
            self.assertIn("file_name", r)
            self.assertIn("meta", r)

    def test_parse_non_existent_file(self):
        with self.assertRaises(Exception) as ctx:
            self.parser.parse_file(os.path.join(self.tmp_dir, "nope.txt"))
        self.assertEqual(ctx.exception.code, "DOCUMENT_NOT_FOUND")

    def test_parse_non_existent_folder(self):
        with self.assertRaises(Exception) as ctx:
            self.parser.parse_folder(os.path.join(self.tmp_dir, "nope_folder"))
        self.assertEqual(ctx.exception.code, "FOLDER_NOT_FOUND")


if __name__ == "__main__":
    unittest.main()
