from __future__ import annotations

import gzip
import json
import os
import re
import tarfile
import tempfile
import zipfile
from html.parser import HTMLParser
from typing import Dict, List, Optional, Tuple

# 第三方解析依赖 — 改为可选 import:
# 这些库各自较重(pandas / PyPDF2 / python-docx / python-pptx),不应该让
# document_parser_module 仅仅"被 import"(如 ABC 守护扫描) 就强制要求安装。
# 真正使用对应解析能力时, 由 _parse_xxx 方法自己检查 None 并报错。
try:
    import xmltodict  # type: ignore
except ImportError:
    xmltodict = None  # type: ignore

try:
    from PyPDF2 import PdfReader  # type: ignore
except ImportError:
    PdfReader = None  # type: ignore

try:
    from docx import Document  # type: ignore
except ImportError:
    Document = None  # type: ignore

try:
    import pandas as pd  # type: ignore
except ImportError:
    pd = None  # type: ignore

try:
    from pptx import Presentation  # type: ignore
except ImportError:
    Presentation = None  # type: ignore

from .base import BaseDocumentParser
from ..config.config import DEFAULT_ENCODING, SUPPORTED_FILE_TYPES
from ..utils.tool_functions import check_file_type, get_file_extension


# ------------------------
# 依赖模块的“兼容导入”
# ------------------------
try:
    from common_utils_module.core.impl import CommonUtils  # type: ignore
except Exception:  # pragma: no cover

    class CommonUtils:  # noqa: D401
        """兼容实现：当系统通用工具模块不可用时，提供最小可用的 clean_text。"""

        @staticmethod
        def clean_text(text: str) -> str:
            text = text.replace("\u00a0", " ")
            text = re.sub(r"\r\n?", "\n", text)
            # 仅折叠连续空格，保留 TAB（用于 csv/excel 表格分隔）
            text = re.sub(r"[ ]+", " ", text)
            text = re.sub(r"\n{3,}", "\n\n", text)
            return text.strip()


try:
    from log_module.core.impl import SystemLogger  # type: ignore
except Exception:  # pragma: no cover
    import logging

    class SystemLogger:  # noqa: D401
        """兼容实现：当系统日志模块不可用时，使用标准 logging。"""

        def __init__(self):
            self._logger = logging.getLogger("document_parser_module")
            if not self._logger.handlers:
                handler = logging.StreamHandler()
                formatter = logging.Formatter(
                    fmt="%(asctime)s %(levelname)s %(name)s - %(message)s"
                )
                handler.setFormatter(formatter)
                self._logger.addHandler(handler)
                self._logger.setLevel(logging.INFO)

        def info(self, msg: str):
            self._logger.info(msg)

        def warning(self, msg: str):
            self._logger.warning(msg)

        def error(self, msg: str):
            self._logger.error(msg)


try:
    from exception_module.core.impl import RAGException  # type: ignore
except Exception:  # pragma: no cover

    class RAGException(Exception):
        """兼容实现：当系统异常模块不可用时，提供最小可用的 RAGException。"""

        def __init__(self, code: str, message: str):
            super().__init__(f"{code}: {message}")
            self.code = code
            self.message = message


# ------------------------
# 压缩包解析参数 (防 zip bomb): 数量/单成员/累计三重上限, 超出只列清单不展开
# ------------------------
ARCHIVE_MAX_MEMBERS = 200
ARCHIVE_MEMBER_MAX_BYTES = 20 * 1024 * 1024
ARCHIVE_MAX_TOTAL_BYTES = 80 * 1024 * 1024
_ARCHIVE_SUFFIXES = (".zip", ".tar", ".tar.gz", ".tgz", ".tar.bz2", ".tar.xz", ".gz")


# ------------------------
# 内部辅助：HTML -> 纯文本
# ------------------------
class _HTMLTextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self._chunks: List[str] = []
        self._skip_depth = 0  # script/style 内的文本不是正文, 不收集

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in ("script", "style"):
            self._skip_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in ("script", "style") and self._skip_depth > 0:
            self._skip_depth -= 1

    def handle_data(self, data: str) -> None:
        if data and self._skip_depth == 0:
            self._chunks.append(data)

    def get_text(self) -> str:
        return "".join(self._chunks)


class LocalDocumentParser(BaseDocumentParser):
    """本地文档解析实现类。

    负责解析 txt/pdf/docx/md/py/excel/ppt/pptx/csv/json/xml/html 为文本，不做存储，返回统一标准结构。
    压缩包 (zip/tar/gz) 当容器展开：成员逐个解析，产出"清单+正文"单文档 (防炸弹三重上限)。
    其余扩展名按内容嗅探：文本文件 (log/yaml/源码等任意后缀) 当纯文本解析，二进制才拒
    (仅 parse_file；parse_folder 批量扫描仍按 supported_file_types 白名单跳过, 防误吞目录杂物)。
    """

    def __init__(self, deps=None):
        # deps: 可选 BasicDeps,未注入时构造一份(向后兼容)
        if deps is not None:
            self.utils = deps.utils
            self.logger = deps.logger
        else:
            try:
                from deps_module import build_basic_deps
                _deps = build_basic_deps()
                self.utils = _deps.utils
                self.logger = _deps.logger
            except Exception:
                # 极端情况下 deps_module 不可用,沿用兼容 import 的实现
                self.utils = CommonUtils()
                self.logger = SystemLogger()
        self.supported_file_types = list(SUPPORTED_FILE_TYPES)

    # ------------------------
    # 各文件类型解析
    # ------------------------
    def _parse_txt(self, file_path: str) -> str:
        with open(file_path, "r", encoding=DEFAULT_ENCODING, errors="ignore") as f:
            return f.read()

    def _parse_pdf(self, file_path: str) -> str:
        reader = PdfReader(file_path)
        texts: List[str] = []
        for i, page in enumerate(reader.pages, start=1):
            try:
                page_text = page.extract_text() or ""
            except Exception:
                page_text = ""
            if page_text.strip():
                texts.append(f"[Page {i}]\n{page_text}")
        return "\n\n".join(texts)

    def _parse_docx(self, file_path: str) -> str:
        doc = Document(file_path)
        paras = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        return "\n".join(paras)

    def _parse_md(self, file_path: str) -> str:
        with open(file_path, "r", encoding=DEFAULT_ENCODING, errors="ignore") as f:
            md_text = f.read()
        # 优先使用 markdown 库；如未安装则降级为“轻量去标记”
        try:
            import markdown  # type: ignore

            html = markdown.markdown(md_text, extensions=["extra", "tables", "fenced_code"])  # type: ignore
            parser = _HTMLTextExtractor()
            parser.feed(html)
            return parser.get_text()
        except Exception:
            # 轻量降级：删除常见 markdown 标记，保留可读文本
            text = md_text
            # code fences
            text = re.sub(r"```[\s\S]*?```", lambda m: m.group(0).strip("`") , text)
            # inline code/backticks
            text = re.sub(r"`([^`]+)`", r"\1", text)
            # headings
            text = re.sub(r"^#{1,6}\s*", "", text, flags=re.MULTILINE)
            # emphasis/bold
            text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
            text = re.sub(r"\*([^*]+)\*", r"\1", text)
            text = re.sub(r"_([^_]+)_", r"\1", text)
            # links [text](url) -> text
            text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1", text)
            return text

    def _parse_py(self, file_path: str) -> str:
        # 保留代码结构与注释：直接读取文本
        with open(file_path, "r", encoding=DEFAULT_ENCODING, errors="ignore") as f:
            return f.read()

    def _parse_html(self, file_path: str) -> str:
        # 复用 _HTMLTextExtractor (跟 _parse_md 的 markdown→HTML 路径同一套),
        # script/style 内容在 extractor 里跳过
        with open(file_path, "r", encoding=DEFAULT_ENCODING, errors="ignore") as f:
            html_text = f.read()
        parser = _HTMLTextExtractor()
        parser.feed(html_text)
        return parser.get_text()

    @staticmethod
    def _sniff_is_text(file_path: str) -> bool:
        """文本/二进制嗅探: 头 8KB 含 NUL 字节判二进制 (git 同款启发式)。
        常见二进制 (图片/压缩包/可执行) 头部都有 NUL; 文本编码 (utf-8/gbk/latin-1)
        正文不会出现 NUL。空文件按文本放行 (上游按空内容 skip)。"""
        with open(file_path, "rb") as f:
            head = f.read(8192)
        return b"\x00" not in head

    # ------------------------
    # 压缩包 (zip/tar/gz): 容器不是"不可解析的二进制" — 展开成员逐个解析
    # ------------------------
    @staticmethod
    def _is_archive(file_name: str) -> bool:
        return str(file_name or "").lower().endswith(_ARCHIVE_SUFFIXES)

    @staticmethod
    def _fix_zip_name(info: "zipfile.ZipInfo") -> str:
        """zip 规范: flag 0x800 表示文件名 UTF-8; 未设时 zipfile 按 cp437 解码 —
        Windows 中文压缩包实际是 GBK, 还原一次, 否则清单里全是乱码。"""
        name = info.filename
        if not (info.flag_bits & 0x800):
            try:
                name = name.encode("cp437").decode("gbk")
            except Exception:
                pass
        return name

    def _read_zip_members(self, file_path: str) -> List[Tuple[str, int, Optional[bytes], str]]:
        out: List[Tuple[str, int, Optional[bytes], str]] = []
        total = 0
        with zipfile.ZipFile(file_path) as zf:
            infos = [i for i in zf.infolist() if not i.is_dir()]
            for idx, info in enumerate(infos):
                if idx >= ARCHIVE_MAX_MEMBERS:
                    out.append((f"(其余 {len(infos) - idx} 个成员)", 0, None, "超出成员数上限, 仅列出"))
                    break
                mname = self._fix_zip_name(info)
                size = int(info.file_size or 0)
                if info.flag_bits & 0x1:
                    out.append((mname, size, None, "加密成员, 跳过"))
                    continue
                if size > ARCHIVE_MEMBER_MAX_BYTES:
                    out.append((mname, size, None, "成员过大, 跳过"))
                    continue
                if total + size > ARCHIVE_MAX_TOTAL_BYTES:
                    out.append((mname, size, None, "解压总量达上限, 跳过"))
                    continue
                try:
                    data = zf.read(info)
                except Exception as e:
                    out.append((mname, size, None, f"读取失败: {e}"))
                    continue
                total += len(data)
                out.append((mname, size, data, ""))
        return out

    def _read_tar_members(self, file_path: str) -> List[Tuple[str, int, Optional[bytes], str]]:
        out: List[Tuple[str, int, Optional[bytes], str]] = []
        total = 0
        count = 0
        with tarfile.open(file_path, "r:*") as tf:
            for m in tf:
                if not m.isfile():
                    continue
                if count >= ARCHIVE_MAX_MEMBERS:
                    out.append(("(更多成员)", 0, None, "超出成员数上限, 仅列出"))
                    break
                count += 1
                size = int(m.size or 0)
                if size > ARCHIVE_MEMBER_MAX_BYTES:
                    out.append((m.name, size, None, "成员过大, 跳过"))
                    continue
                if total + size > ARCHIVE_MAX_TOTAL_BYTES:
                    out.append((m.name, size, None, "解压总量达上限, 跳过"))
                    continue
                f = tf.extractfile(m)
                if f is None:
                    out.append((m.name, size, None, "无法读取"))
                    continue
                data = f.read()
                total += len(data)
                out.append((m.name, size, data, ""))
        return out

    def _read_gz_member(self, file_path: str) -> List[Tuple[str, int, Optional[bytes], str]]:
        """裸 .gz 单文件 (如 app.log.gz)。.tar.gz 走 tar 分支, 不到这里。"""
        inner = os.path.basename(file_path)
        if inner.lower().endswith(".gz"):
            inner = inner[:-3] or inner
        with gzip.open(file_path, "rb") as f:
            data = f.read(ARCHIVE_MAX_TOTAL_BYTES + 1)
        if len(data) > ARCHIVE_MAX_TOTAL_BYTES:
            return [(inner, len(data), None, "解压超上限, 跳过")]
        return [(inner, len(data), data, "")]

    def _member_to_text(self, member_name: str, data: bytes) -> Tuple[str, str]:
        """单个成员 → (正文, 跳过原因)。已知格式落临时文件复用 parse_file;
        其余按内容嗅探 (与顶层文件同一套规则)。"""
        if self._is_archive(member_name):
            return "", "嵌套压缩包, 不展开"
        suffix = os.path.splitext(member_name)[1].lower()
        if suffix in self.supported_file_types:
            fd, tmp = tempfile.mkstemp(suffix=suffix or ".txt")
            try:
                with os.fdopen(fd, "wb") as f:
                    f.write(data)
                try:
                    return str(self.parse_file(tmp).get("content") or ""), ""
                except Exception as e:
                    return "", f"成员解析失败: {getattr(e, 'message', e)}"
            finally:
                try:
                    os.unlink(tmp)
                except OSError:
                    pass
        if b"\x00" in data[:8192]:
            return "", "二进制, 仅列清单"
        return data.decode(DEFAULT_ENCODING, errors="ignore"), ""

    def _parse_archive(self, file_path: str) -> str:
        """压缩包 → "清单 + 各成员正文" 单文档。全二进制成员的包也产出清单
        (Agent 至少能回答"包里有什么")。"""
        name = os.path.basename(file_path)
        low = name.lower()
        try:
            if low.endswith(".zip"):
                members = self._read_zip_members(file_path)
            elif low.endswith((".tar", ".tar.gz", ".tgz", ".tar.bz2", ".tar.xz")):
                members = self._read_tar_members(file_path)
            else:
                members = self._read_gz_member(file_path)
        except RAGException:
            raise
        except Exception as e:
            raise RAGException("DOCUMENT_PARSE_FAILED", f"压缩包解析失败：{file_path}，原因：{e}")

        manifest: List[str] = []
        bodies: List[str] = []
        extracted = 0
        for mname, size, data, note in members:
            line = f"- {mname} ({size} bytes)"
            if data is not None:
                text, skip_reason = self._member_to_text(mname, data)
                if text.strip():
                    extracted += 1
                    bodies.append(f"=== {mname} ===\n{text}")
                elif skip_reason:
                    note = skip_reason
            if note:
                line += f" — {note}"
            manifest.append(line)

        parts = [
            f"[压缩包 {name} — {len(members)} 个成员, 提取正文 {extracted} 个]",
            "清单:",
            "\n".join(manifest) if manifest else "(空压缩包)",
        ]
        if bodies:
            parts.append("\n\n".join(bodies))
        return "\n".join(parts)

    def _parse_excel(self, file_path: str) -> str:
        # 以工作表组织文本：SheetName + 表格内容
        try:
            sheets = pd.read_excel(file_path, sheet_name=None, dtype=str)  # type: ignore
        except Exception as e:
            raise RAGException("EXCEL_PARSE_ERROR", f"excel文件解析异常：{file_path}，原因：{e}")

        parts: List[str] = []
        for sheet_name, df in sheets.items():
            df = df.fillna("")
            table_text = df.to_csv(sep="\t", index=False)
            parts.append(f"[Sheet: {sheet_name}]\n{table_text}")
        return "\n\n".join(parts)

    def _parse_ppt(self, file_path: str) -> str:
        try:
            pres = Presentation(file_path)
        except Exception as e:
            raise RAGException("PPT_PARSE_ERROR", f"ppt文件解析异常：{file_path}，原因：{e}")

        parts: List[str] = []
        for idx, slide in enumerate(pres.slides, start=1):
            slide_texts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = (shape.text or "").strip()
                    if t:
                        slide_texts.append(t)
            if slide_texts:
                parts.append(f"[Slide {idx}]\n" + "\n".join(slide_texts))
            else:
                parts.append(f"[Slide {idx}]\n")
        return "\n\n".join(parts)

    def _parse_csv(self, file_path: str) -> str:
        try:
            df = pd.read_csv(file_path, dtype=str)
        except Exception as e:
            raise RAGException("DOCUMENT_PARSE_FAILED", f"csv文件解析异常：{file_path}，原因：{e}")
        df = df.fillna("")
        return df.to_csv(sep="\t", index=False)

    def _parse_json(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding=DEFAULT_ENCODING, errors="ignore") as f:
                obj = json.load(f)
        except Exception as e:
            raise RAGException("JSON_PARSE_ERROR", f"json文件解析异常：{file_path}，原因：{e}")
        return json.dumps(obj, ensure_ascii=False, indent=2)

    def _parse_xml(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding=DEFAULT_ENCODING, errors="ignore") as f:
                xml_text = f.read()
            obj = xmltodict.parse(xml_text)
        except Exception as e:
            raise RAGException("XML_PARSE_ERROR", f"xml文件解析异常：{file_path}，原因：{e}")
        # 统一输出为“可读文本”：使用 JSON 形式展示结构
        return json.dumps(obj, ensure_ascii=False, indent=2)

    # ------------------------
    # 抽象方法实现
    # ------------------------
    def parse_file(self, file_path: str) -> Dict:
        if not os.path.isfile(file_path):
            raise RAGException("DOCUMENT_NOT_FOUND", f"文档文件不存在：{file_path}")

        ext = get_file_extension(file_path)
        file_name = os.path.basename(file_path)

        try:
            if ext == ".txt":
                content = self._parse_txt(file_path)
            elif ext == ".pdf":
                content = self._parse_pdf(file_path)
            elif ext == ".docx":
                content = self._parse_docx(file_path)
            elif ext == ".md":
                content = self._parse_md(file_path)
            elif ext == ".py":
                content = self._parse_py(file_path)
            elif ext in (".xlsx", ".xls"):
                content = self._parse_excel(file_path)
            elif ext in (".ppt", ".pptx"):
                content = self._parse_ppt(file_path)
            elif ext == ".csv":
                content = self._parse_csv(file_path)
            elif ext == ".json":
                content = self._parse_json(file_path)
            elif ext == ".xml":
                content = self._parse_xml(file_path)
            elif ext in (".html", ".htm"):
                content = self._parse_html(file_path)
            elif self._is_archive(file_name):
                # 压缩包在二进制嗅探之前接住 (zip/tar 头部含 NUL, 嗅探会误拒)
                content = self._parse_archive(file_path)
            else:
                # 未知扩展名: 按内容嗅探, 不按后缀清单拒 — 后缀清单是打地鼠
                # (.log/.yaml/.sql/各语言源码无穷尽)。文本就当纯文本解析, 二进制才拒。
                if not self._sniff_is_text(file_path):
                    raise RAGException(
                        "UNSUPPORTED_FILE_TYPE",
                        f"二进制文件无法提取文本：{ext or '(无扩展名)'}（文件：{file_path}）",
                    )
                content = self._parse_txt(file_path)

            # 基础清洗（优先使用系统 CommonUtils.clean_text）
            if hasattr(self.utils, "clean_text"):
                content = self.utils.clean_text(content)
            else:  # pragma: no cover
                content = str(content).strip()

            self.logger.info(f"解析成功：{file_path}")
            return {"content": content, "file_name": file_name, "meta": {"ext": ext}}

        except RAGException:
            self.logger.error(f"解析失败：{file_path}")
            raise
        except Exception as e:
            self.logger.error(f"解析失败：{file_path}，原因：{e}")
            raise RAGException("DOCUMENT_PARSE_FAILED", f"文档解析失败：{file_path}，原因：{e}")

    def parse_folder(self, folder_path: str) -> List[Dict]:
        if not os.path.isdir(folder_path):
            raise RAGException("FOLDER_NOT_FOUND", f"文件夹不存在：{folder_path}")

        results: List[Dict] = []
        for root, _, files in os.walk(folder_path):
            for name in files:
                fp = os.path.join(root, name)
                if not check_file_type(fp, self.supported_file_types):
                    self.logger.warning(f"跳过不支持的文件：{fp}")
                    continue
                results.append(self.parse_file(fp))

        self.logger.info(f"批量解析完成：{folder_path}，共 {len(results)} 个文件")
        return results
